from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
import math
import pandas as pd
import os
import subprocess
from pdf2image import convert_from_path

CM = 1/361507
DISTANCE_FROM_CENTER_CM = 8.5
HORIZONTAL_CENTER = 10 * 361507
VERTICAL_CENTER = 10 * 361507

N = 5 # グラフの頂点数
CENTRAL_ANGLE = 2 * math.pi / N
ANGLE_DEGREE_90 = 1/2 * math.pi
ANGLE_DEGREE_18 = 18 * math.pi / 180

# path
LIBRE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
PPT_TEMPLETE_PATH = "../../const/graph_openwork.pptm"
OUTPUT_PARENT_DIR_PATH = "../../output/"
INPUT_DIR_PATH = "../../output_openwork/"

PRS = Presentation(PPT_TEMPLETE_PATH)

def PPT2PDF(output_path):
    if(os.path.exists(PPT_TEMPLETE_PATH)):
        if not (os.path.exists(output_path)):
            os.mkdir(output_path)
        
        command = f"{LIBRE_PATH} --headless --nologo --nofirststartwizard --convert-to pdf --outdir {output_path} {PPT_TEMPLETE_PATH}".split(
                " ")
        subprocess.run(command)

def PDF2PNG(output_path):
    pdf_path = output_path + '/graph_openwork.pdf'
    # pdfから画像に変換
    page = convert_from_path(str(pdf_path), dpi=200, fmt='png')[0]
    page.save(output_path + '/graph_openwork.png', "PNG")

def delete_PDF(output_path):
    pdf_path = output_path + '/graph_openwork.pdf'
    os.remove(pdf_path)

def calc_score(score):
    try:
        score = float(score)
    except Exception:
        return 0
    score = score * 2 * 10  #100点換算
    score = score + 0.3 * (100-score) #平均化
    score = round(score, 1)

    return score

def record_data(output_path, openwork_label_list, openwork_data_list):
    dic_openwork_data = dict(zip(openwork_label_list, openwork_data_list))

    df_openwork = pd.DataFrame(index=[], columns=openwork_label_list)
    df_openwork = df_openwork.append(dic_openwork_data, ignore_index=True)

    df_openwork.to_csv(output_path + '/openwork_data.csv', header=openwork_label_list, encoding="utf-8", index=False)

def calc_coorfinate(score, vertex_num):
    left = round(math.cos(CENTRAL_ANGLE * vertex_num - ANGLE_DEGREE_90) * (8.5 * score/100) / CM + HORIZONTAL_CENTER)
    top = round(math.sin(CENTRAL_ANGLE * vertex_num - ANGLE_DEGREE_90) * (8.5 * score/100) / CM + VERTICAL_CENTER)

    dic_coorfinate = {
        'left': left,
        'top': top,
    }

    return dic_coorfinate

def chart_function(graph, after_graph):
    graph.chart.replace_data(after_graph)

def dot_function(shape, score, coorfinate_num):
    if(coorfinate_num >= 5):
        return False
    dic_coorfinate = calc_coorfinate(score, coorfinate_num)
    shape.left = dic_coorfinate['left']
    shape.top = dic_coorfinate['top']

    return True

def create_graph(c_data, score_tuple):
    shapes = PRS.slides[0].shapes

    coorfinate_num = 0
    for shape in shapes:
        if(shape.has_chart):
            chart_function(shape, c_data)
        
        if(shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE):  
            if(dot_function(shape, score_tuple[coorfinate_num], coorfinate_num)):
                coorfinate_num += 1
            else:
                continue
    
    PRS.save(PPT_TEMPLETE_PATH)

# 企業リストを読み込み
col_names = ['待遇面の満足度', '社員の士気', '風通しの良さ', '社員の相互尊重', '20代成長環境', '人材の長期育成', '法令遵守意識', '人事評価の適正感']

# 業界毎のcsv一覧を取得
indestry_csv_list = os.listdir(INPUT_DIR_PATH)

for industry_csv in indestry_csv_list:
    industry_name = industry_csv.replace('.csv', '')
    print(industry_name + ': START!!')

    if(industry_name == '.DS_Store'):
        continue

    # 企業リストを読み込み  
    input_filename = INPUT_DIR_PATH + industry_csv
    df_company = pd.read_csv(input_filename, encoding="utf-8")

    for index, company in df_company.iterrows():
        company_name = company['企業名'].replace(' ', '')
        print(company_name)
        output_path = OUTPUT_PARENT_DIR_PATH + industry_name + '/' + str(company['証券コード']) + '_' + company_name

        if(os.path.exists(output_path + '/graph_openwork.png')):
            continue

        score_tuple = (
            calc_score(company['待遇面の満足度']),
            calc_score(company['20代成長環境']) * 2/3 + calc_score(company['人材の長期育成']) * 1/3,
            calc_score(company['社員の士気']),
            calc_score(company['人事評価の適正感']),
            calc_score(company['風通しの良さ'])/2 + calc_score(company['社員の相互尊重'])/4 + calc_score(company['法令遵守意識'])/4,
        )

        graph_label_list = ['給与・待遇', '成長環境', 'やりがい', '人事評価', '職場環境']
        c_data = CategoryChartData()
        c_data.categories = graph_label_list
        c_data.add_series('X', score_tuple)

        create_graph(c_data, score_tuple)

        PPT2PDF(output_path)
        PDF2PNG(output_path)
        delete_PDF(output_path)
        record_data(output_path, graph_label_list, list(score_tuple))