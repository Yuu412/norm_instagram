from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
import math
import pandas as pd
import os
import subprocess
from pdf2image import convert_from_path

# path
LIBRE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
PPT_TEMPLETE_PATH = "../../const/stability_profitability_growth.pptm"
INPUT_PATH = '../../input'
OUTPUT_PATH = "../../output/"

PRS = Presentation(PPT_TEMPLETE_PATH)

LABEL_LISTS = [
    ['安定性_1', '安定性_2', '安定性_3'],
    ['成長性_1', '成長性_2', '成長性_3'],
    ['収益性_1', '収益性_2', '収益性_3'],
]
GRAPH_LABEL_LIST = ['2020', '2021', '2022']

def PPT2PDF(output_path):
    if(os.path.exists(PPT_TEMPLETE_PATH)):
        if not (os.path.exists(output_path)):
            os.mkdir(output_path)
        
        command = f"{LIBRE_PATH} --headless --nologo --nofirststartwizard --convert-to pdf --outdir {output_path} {PPT_TEMPLETE_PATH}".split(
                " ")
        subprocess.run(command)

def PDF2PNG(output_path):
    pdf_path = output_path + '/stability_profitability_growth.pdf'
    # pdfから画像に変換
    page = convert_from_path(str(pdf_path), dpi=200, fmt='png')[0]
    page.save(output_path + '/stability_profitability_growth.png', "PNG")

def delete_PDF(output_path):
    pdf_path = output_path + '/stability_profitability_growth.pdf'
    os.remove(pdf_path)

def calc_y_range(taple):
    y_max = math.ceil(max(taple)/10) * 10
    y_min = math.floor(min(taple)/10) * 10

    if(y_min < 0):
        y_min = 0

    if(y_max > 100):
        y_max = 100

    dic_range = {
        'max': y_max,
        'min': y_min
    }
    return dic_range


def replace_graph_data(c_data, score_tuple):
    # テンプレートから要素の配列を取得
    shapes = PRS.slides[0].shapes
    
    #グラフの要素を取得
    graph_shape = shapes[0]
    
    # グラフのデータを置き換え
    graph_shape.chart.replace_data(c_data)

    # グラフのY軸の最大値・最小値を設定
    dic_range = calc_y_range(score_tuple)
    graph_shape.chart.value_axis.maximum_scale = dic_range['max']
    graph_shape.chart.value_axis.minimum_scale = dic_range['min']

    # テンプレートの保存
    PRS.save(PPT_TEMPLETE_PATH)

    return True

# 業界名のcsv一覧を取得
indestry_csv_list = os.listdir(INPUT_PATH)
for industry_csv in indestry_csv_list:
    industry = industry_csv.replace('.csv', '')
    print(industry + ': START!!')

    if(industry == '.DS_Store'):
        continue

    # 企業リストを読み込み  
    input_filename = OUTPUT_PATH + industry + '/company_score.csv'
    df_company_score = pd.read_csv(input_filename, encoding="utf-8")

    for index, company in df_company_score.iterrows():
        company_name = company['企業名'].replace(' ', '')
        print(company_name)
    
        output_path = OUTPUT_PATH + industry + '/' + str(company['証券コード']) + '_' + company_name

        for label_list in LABEL_LISTS:
            label_name = label_list[0].split('_')[0]
            #if(os.path.exists(output_path + '/' + label_name + '.png')):
                #continue

            score_tuple = (
                company[label_list[0]],
                company[label_list[1]],
                company[label_list[2]],
            )
            
            c_data = CategoryChartData()
            c_data.categories = GRAPH_LABEL_LIST
            c_data.add_series('X', score_tuple)

            replace_graph_data(c_data, score_tuple)

            PPT2PDF(output_path)
            PDF2PNG(output_path)
            delete_PDF(output_path)
            os.rename(
                output_path + '/stability_profitability_growth.png',
                output_path + '/' + label_name + '.png'
            )