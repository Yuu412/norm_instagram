from cv2 import floodFill
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
import math
import pandas as pd
import os
import subprocess
from pdf2image import convert_from_path
from pptx.enum.chart import XL_AXIS_CROSSES

# graph types
COLUMN_CLUSTERED = 51
LINE_MARKERS = 65

# path
LIBRE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
PPT_TEMPLETE_PATH = "../../const/graph_performance_templete.pptm"
OUTPUT_PARENT_DIR_PATH = "../../output/"
INPUT_DIR_PATH = "../../output_performance/"

PRS = Presentation(PPT_TEMPLETE_PATH)

def PPT2PDF(output_path):
    if(os.path.exists(PPT_TEMPLETE_PATH)):
        if not (os.path.exists(output_path)):
            os.mkdir(output_path)
        
        command = f"{LIBRE_PATH} --headless --nologo --nofirststartwizard --convert-to pdf --outdir {output_path} {PPT_TEMPLETE_PATH}".split(" ")
        subprocess.run(command)

def PDF2PNG(output_path):
    pdf_path = output_path + '/graph_performance_templete.pdf'
    # pdfから画像に変換
    page = convert_from_path(str(pdf_path), dpi=200, fmt='png')[0]
    page.save(output_path + '/graph_performance_templete.png', 'PNG')

def delete_PDF(output_path):
    pdf_path = output_path + '/graph_performance_templete.pdf'

    os.remove(pdf_path)

def ret_year(year):
    try:
        int_year = int(year)
    except Exception:
        int_year = '-'
    return int_year

def arange_value(value):
    if(type(value) is int):
        return round(value)
    
    if(value == '単' or value == '連' or value == '-' or value == '---'):
        return -1

    return int(value.split('.')[0].replace(',',''))

def arange_taple(taple_value):
    while(max(taple_value) > 20000):
        taple_value = tuple([round(value / 10) for value in taple_value])

    return taple_value

def record_data(output_path, index_list, column_list, sales_list, profit_list):
    df_performance = pd.DataFrame(list(zip(sales_list, profit_list)), index=index_list, columns=column_list)
    pdf_path = output_path + '/performance_data.csv'

    df_performance.to_csv(pdf_path, encoding="utf-8", index=index_list)

def calc_y_range(taple):
    y_max = round(max(taple) * 1.2, -3)
    y_min = round(min(taple) / 1.2, -3)

    if(y_min <= 2000):
        y_min = 0
    
    if((y_max - y_min) < 2000):
        if(y_min == 0):
            y_max = y_max + 1000
        else:
            y_min = y_min - 1000

    dic_range = {
        'max': y_max,
        'min': y_min
    }
    return dic_range


def create_graph(chart_data, graph_type_num, value_taple):
    shapes = PRS.slides[0].shapes

    for shape in shapes:
        if(shape.chart.chart_type == graph_type_num):
            shape.chart.replace_data(chart_data)
            
            # グラフのY軸の最大値・最小値を設定
            dic_range = calc_y_range(value_taple)
            shape.chart.value_axis.maximum_scale = dic_range['max']
            shape.chart.value_axis.minimum_scale = dic_range['min']
    
    PRS.save(PPT_TEMPLETE_PATH)

# 業界毎のcsv一覧を取得
indestry_csv_list = os.listdir(INPUT_DIR_PATH)

for industry_csv in indestry_csv_list:
    industry_name = industry_csv.replace('.csv', '')
    print(industry_name + ': START!!')

    if(industry_name == '.DS_Store'):
        continue

    # すでにフォルダが作成済みの場合
    if not os.path.exists(OUTPUT_PARENT_DIR_PATH + industry_name):
        os.mkdir(OUTPUT_PARENT_DIR_PATH + industry_name)

    # 企業リストを読み込み
    input_filename = INPUT_DIR_PATH + industry_csv
    df_company = pd.read_csv(input_filename, encoding="utf-8")

    for index, company in df_company.iterrows():
        company_name = company['企業名'].replace(' ', '')
        foldername = str(company['証券コード']) + '_' + company_name
        if(os.path.exists(OUTPUT_PARENT_DIR_PATH + industry_name + '/' + foldername)):
            continue
        print(company_name)
        year_list = [
            ret_year(company['prior4Year']),
            ret_year(company['prior3Year']),
            ret_year(company['prior2Year']),
            ret_year(company['prior1Year']),
        ]

        sales_taple = (
            arange_value(company['prior4YearSales']),
            arange_value(company['prior3YearSales']),
            arange_value(company['prior2YearSales']),
            arange_value(company['prior1YearSales']),
        )
        sales_taple = arange_taple(sales_taple)

        profit_taple = (
            arange_value(company['prior4YearProfit']),
            arange_value(company['prior3YearProfit']),
            arange_value(company['prior2YearProfit']),
            arange_value(company['prior1YearProfit']),
        )
        profit_taple = arange_taple(profit_taple)

        sales_data = CategoryChartData()
        sales_data.categories = year_list
        sales_data.add_series('売上高', sales_taple)

        profit_data = CategoryChartData()
        profit_data.categories = year_list
        profit_data.add_series('営業利益', profit_taple)

        try:
            create_graph(sales_data, COLUMN_CLUSTERED, sales_taple)
            create_graph(profit_data, LINE_MARKERS, profit_taple)
        except Exception:
            continue

        output_path = OUTPUT_PARENT_DIR_PATH + industry_name + '/' + foldername

        PPT2PDF(output_path)
        PDF2PNG(output_path)
        delete_PDF(output_path)
        record_data(output_path, year_list, ['売上高', '営業利益'], list(sales_taple), list(profit_taple))